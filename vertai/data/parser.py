"""Document parser module supporting multiple file formats."""

from __future__ import annotations

import re
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Any, Optional, TYPE_CHECKING

if TYPE_CHECKING:
    # Type-only imports for the optional doc-parser extras. The actual
    # imports happen lazily inside ``parse()`` methods guarded by try/except.
    from docx.table import Table as _DocxTable
    from openpyxl.worksheet.worksheet import Worksheet as _OpenpyxlWorksheet


class BaseParser(ABC):
    """Abstract base parser defining the parsing interface."""

    @abstractmethod
    def parse(self, file_path: Path) -> dict[str, Any]:
        """
        Parse document and return structured content.

        Args:
            file_path: Path to the document file.

        Returns:
            Dictionary with keys:
                - text: Full document text as string
                - metadata: Document metadata dictionary
                - chunks: List of content chunks with position info

        Raises:
            ImportError: If required dependency is not installed.
        """
        pass

    @staticmethod
    def _create_result(
        text: str,
        metadata: dict[str, Any],
        chunks: Optional[list[dict[str, Any]]] = None,
    ) -> dict[str, Any]:
        """Create standardized parsing result."""
        return {
            "text": text,
            "metadata": metadata,
            "chunks": chunks or [],
        }


class PDFParser(BaseParser):
    """Parse PDF documents using PyMuPDF."""

    def parse(self, file_path: Path) -> dict[str, Any]:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError(
                "PyMuPDF is required for PDF parsing. "
                "Install with: pip install PyMuPDF"
            )

        doc = fitz.open(str(file_path))
        try:
            text_parts = []
            chunks = []
            metadata = {
                "pages": len(doc),
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "subject": doc.metadata.get("subject", ""),
                "creator": doc.metadata.get("creator", ""),
                "producer": doc.metadata.get("producer", ""),
                "creation_date": doc.metadata.get("creationDate", ""),
                "format": "PDF",
            }

            for page_num, page in enumerate(doc, start=1):
                page_text = page.get_text()
                text_parts.append(page_text)
                chunks.append({
                    "text": page_text.strip(),
                    "page": page_num,
                })

            return self._create_result(
                text="\n\n".join(text_parts),
                metadata=metadata,
                chunks=chunks,
            )
        finally:
            doc.close()


class WordParser(BaseParser):
    """Parse Word documents using python-docx."""

    def parse(self, file_path: Path) -> dict[str, Any]:
        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx is required for Word parsing. "
                "Install with: pip install python-docx"
            )

        doc = Document(str(file_path))
        text_parts = []
        chunks = []

        core_props = doc.core_properties
        metadata = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
            "format": "DOCX",
        }

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
                chunks.append({
                    "text": para.text,
                    "type": "paragraph",
                    "style": para.style.name if para.style else "",
                })

        for table in doc.tables:
            table_text = self._extract_table(table)
            if table_text.strip():
                text_parts.append(table_text)
                chunks.append({
                    "text": table_text,
                    "type": "table",
                })

        return self._create_result(
            text="\n\n".join(text_parts),
            metadata=metadata,
            chunks=chunks,
        )

    def _extract_table(self, table: "_DocxTable") -> str:
        """Extract table content as text."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)


class ExcelParser(BaseParser):
    """Parse Excel documents using openpyxl."""

    def parse(self, file_path: Path) -> dict[str, Any]:
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError(
                "openpyxl is required for Excel parsing. "
                "Install with: pip install openpyxl"
            )

        wb = load_workbook(str(file_path), data_only=True)
        text_parts = []
        chunks = []
        metadata = {
            "sheets": len(wb.sheetnames),
            "sheet_names": wb.sheetnames,
            "format": "XLSX",
        }

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = self._extract_sheet(ws)
            if sheet_text.strip():
                text_parts.append(f"### {sheet_name}\n{sheet_text}")
                chunks.append({
                    "text": sheet_text,
                    "sheet": sheet_name,
                    "rows": ws.max_row,
                    "columns": ws.max_column,
                })

        wb.close()
        return self._create_result(
            text="\n\n".join(text_parts),
            metadata=metadata,
            chunks=chunks,
        )

    def _extract_sheet(self, worksheet: "_OpenpyxlWorksheet") -> str:
        """Extract worksheet content as text."""
        rows = []
        for row in worksheet.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(cells):
                rows.append(" | ".join(cells))
        return "\n".join(rows)


class PPTParser(BaseParser):
    """Parse PowerPoint documents using python-pptx."""

    def parse(self, file_path: Path) -> dict[str, Any]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint parsing. "
                "Install with: pip install python-pptx"
            )

        prs = Presentation(str(file_path))
        text_parts = []
        chunks = []

        core_props = prs.core_properties
        metadata = {
            "slides": len(prs.slides),
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "format": "PPTX",
        }

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = self._extract_slide(slide)
            if slide_text.strip():
                text_parts.append(f"## Slide {slide_num}\n{slide_text}")
                chunks.append({
                    "text": slide_text,
                    "slide": slide_num,
                })

        return self._create_result(
            text="\n\n".join(text_parts),
            metadata=metadata,
            chunks=chunks,
        )

    def _extract_slide(self, slide: Any) -> str:
        """Extract slide content as text."""
        shapes_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shapes_text.append(shape.text)
        return "\n".join(shapes_text)


class MarkdownParser(BaseParser):
    """Parse Markdown documents using built-in parsing."""

    # Regex patterns for markdown elements
    HEADER_PATTERN = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)
    CODE_BLOCK_PATTERN = re.compile(r"```[\s\S]*?```", re.MULTILINE)

    def parse(self, file_path: Path) -> dict[str, Any]:
        file_path = Path(file_path) if isinstance(file_path, str) else file_path
        content = file_path.read_text(encoding="utf-8")

        headers = self.HEADER_PATTERN.findall(content)
        code_blocks = self.CODE_BLOCK_PATTERN.findall(content)

        metadata = {
            "headers": [{"level": len(h), "text": t} for h, t in headers],
            "header_count": len(headers),
            "code_blocks": len(code_blocks),
            "characters": len(content),
            "lines": content.count("\n") + 1,
            "format": "MD",
        }

        chunks = self._create_chunks(content)

        return self._create_result(
            text=content,
            metadata=metadata,
            chunks=chunks,
        )

    def _create_chunks(self, content: str) -> list[dict[str, Any]]:
        """Create chunks based on markdown sections."""
        chunks = []
        lines = content.split("\n")
        current_chunk: list[str] = []
        current_header = "Introduction"

        for line in lines:
            header_match = self.HEADER_PATTERN.match(line)
            if header_match:
                if current_chunk:
                    chunk_text = "\n".join(current_chunk).strip()
                    if chunk_text:
                        chunks.append({
                            "text": chunk_text,
                            "section": current_header,
                        })
                current_header = header_match.group(2)
                current_chunk = [line]
            else:
                current_chunk.append(line)

        if current_chunk:
            chunk_text = "\n".join(current_chunk).strip()
            if chunk_text:
                chunks.append({
                    "text": chunk_text,
                    "section": current_header,
                })

        return chunks


class DocParser:
    """
    Unified document parser supporting multiple file formats.

    Supported formats: PDF, DOCX, XLSX, PPTX, MD

    Usage:
        >>> from vertai.data import DocParser
        >>> parser = DocParser()
        >>> doc = parser.parse("report.pdf")
        >>> print(doc["text"])
        >>> print(doc["metadata"])
    """

    PARSERS: dict[str, type[BaseParser]] = {
        ".pdf": PDFParser,
        ".docx": WordParser,
        ".xlsx": ExcelParser,
        ".pptx": PPTParser,
        ".md": MarkdownParser,
        ".markdown": MarkdownParser,
    }

    def __init__(self) -> None:
        self._parsers: dict[str, BaseParser] = {}

    def _get_parser(self, extension: str) -> BaseParser:
        """Get or create parser for given extension."""
        ext = extension.lower()
        if ext not in self._parsers:
            parser_class = self.PARSERS.get(ext)
            if not parser_class:
                raise ValueError(
                    f"Unsupported file format: {ext}. "
                    f"Supported formats: {', '.join(self.PARSERS.keys())}"
                )
            self._parsers[ext] = parser_class()
        return self._parsers[ext]

    def parse(self, file_path: str | Path) -> dict[str, Any]:
        """
        Parse a document and return structured content.

        Args:
            file_path: Path to the document file.

        Returns:
            Dictionary containing:
                - text: Full document text
                - metadata: Document metadata (pages, author, etc.)
                - chunks: List of content chunks with position info

        Raises:
            FileNotFoundError: If file doesn't exist.
            ValueError: If file format is not supported.
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        if not path.is_file():
            raise ValueError(f"Not a file: {path}")

        parser = self._get_parser(path.suffix)
        result = parser.parse(path)
        result["metadata"]["file_path"] = str(path.absolute())
        result["metadata"]["file_name"] = path.name

        return result

    def supports(self, file_path: str | Path) -> bool:
        """Check if file format is supported."""
        path = Path(file_path)
        return path.suffix.lower() in self.PARSERS

    @property
    def supported_formats(self) -> list[str]:
        """Return list of supported file extensions."""
        return list(self.PARSERS.keys())
